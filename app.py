import os
import tkinter as tk
from tkinter import filedialog, messagebox
import openpyxl
from docx import Document
from pptx import Presentation
import datetime

def ler_regras_excel(caminho_excel):
    """
    Lê as regras de substituição do arquivo Excel.
    Coluna A = de (Tag), Coluna B = para (Novo Texto).
    Detecta formatações nativas do Excel (Porcentagem, Moeda e Data) e as converte para texto no padrão brasileiro.
    """
    wb = openpyxl.load_workbook(caminho_excel, data_only=True)
    planilha = wb.active
    regras = []
    
    # A iteração agora acessa os objetos da célula (cell), não apenas os valores brutos
    for linha in planilha.iter_rows(min_row=1):
        celula_de = linha[0]
        celula_para = linha[1]
        
        de_texto = celula_de.value
        valor_para = celula_para.value
        
        if de_texto is not None and str(de_texto).strip() != "":
            de_formatado = str(de_texto).strip()
            
            # Formatação do valor da Coluna B
            if valor_para is None:
                para_formatado = ""
                
            # Tratamento para Datas
            elif isinstance(valor_para, datetime.datetime):
                # Extrai o formato do Excel para verificar se tem apenas mês/ano ou data completa
                formato_excel = str(celula_para.number_format).lower()
                if 'mmm' in formato_excel or 'mmm-yy' in formato_excel:
                    # Exemplo: jul/2026
                    para_formatado = valor_para.strftime("%b/%Y").lower()
                else:
                    # Exemplo: 01/07/2026
                    para_formatado = valor_para.strftime("%d/%m/%Y")
                    
            # Tratamento para Números (Moeda, Porcentagem e Comum)
            elif isinstance(valor_para, (int, float)):
                formato_excel = str(celula_para.number_format).lower()
                
                # Verifica se é Porcentagem
                if '%' in formato_excel:
                    # Multiplica por 100 e formata com 2 casas decimais, trocando ponto por vírgula
                    para_formatado = f"{valor_para * 100:.2f}%".replace('.', ',')
                    
                # Verifica se é Moeda/Contábil (busca por $, R$ ou padrões de formatação contábil)
                elif 'r$' in formato_excel or '$' in formato_excel or '_-' in formato_excel:
                    # Formata com separador de milhar e decimal
                    moeda_str = f"{valor_para:,.2f}"
                    # Inverte pontos e vírgulas para o padrão brasileiro (R$ 1.000,00)
                    moeda_str = moeda_str.replace(',', 'X').replace('.', ',').replace('X', '.')
                    para_formatado = f"R$ {moeda_str}"
                    
                else:
                    # Número comum, converte para string
                    para_formatado = str(valor_para).strip()
            else:
                para_formatado = str(valor_para).strip()
                
            regras.append({'de': de_formatado, 'para': para_formatado})
            
    return regras

def processar_word(caminho_entrada, caminho_saida, regras):
    """Aplica as substituições em documentos Word (.docx) iterando sobre parágrafos e tabelas."""
    doc = Document(caminho_entrada)
    
    def substituir_texto(paragraphs):
        for p in paragraphs:
            for regra in regras:
                de, para = regra['de'], regra['para']
                if de in p.text:
                    # Tenta substituir nos 'runs' para preservar formatação (negrito, cor, etc.)
                    for run in p.runs:
                        if de in run.text:
                            run.text = run.text.replace(de, para)
                    # Fallback: Se a tag foi dividida em múltiplos runs pelo Word,
                    # substitui no parágrafo inteiro (pode perder formatações específicas)
                    if de in p.text:
                        p.text = p.text.replace(de, para)

    # Varre parágrafos normais
    substituir_texto(doc.paragraphs)
    
    # Varre tabelas
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                substituir_texto(cell.paragraphs)
                
    doc.save(caminho_saida)

def processar_powerpoint(caminho_entrada, caminho_saida, regras):
    """Aplica as substituições em apresentações PowerPoint (.pptx), incluindo Slide Mestre."""
    prs = Presentation(caminho_entrada)
    
    # Função auxiliar para varrer as formas (shapes) e evitar repetição de código
    def substituir_em_shapes(shapes):
        for shape in shapes:
            # Substituição em caixas de texto padrão
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        for regra in regras:
                            if regra['de'] in run.text:
                                run.text = run.text.replace(regra['de'], regra['para'])
            
            # Substituição dentro de tabelas
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            for run in p.runs:
                                for regra in regras:
                                    if regra['de'] in run.text:
                                        run.text = run.text.replace(regra['de'], regra['para'])

    # 1. Varre os Slides Normais
    for slide in prs.slides:
        substituir_em_shapes(slide.shapes)
        
    # 2. Varre todos os Slides Mestres e seus respectivos Layouts
    for master in prs.slide_masters:
        # Substitui diretamente no Slide Mestre
        substituir_em_shapes(master.shapes)
        
        # Substitui nos Layouts derivados deste Slide Mestre
        for layout in master.slide_layouts:
            substituir_em_shapes(layout.shapes)
                                        
    prs.save(caminho_saida)

def iniciar_automacao():
    """Inicia a interface de seleção e executa a lógica de automação."""
    # Instancia e oculta a janela principal do Tkinter para mostrar apenas o popup
    root = tk.Tk()
    root.withdraw()
    
    # 1. Solicita a seleção do arquivo de relatório
    caminho_arquivo = filedialog.askopenfilename(
        title="Selecione o arquivo Word ou PowerPoint para alterar",
        filetypes=[
            ("Documentos suportados", "*.docx *.pptx"), 
            ("Word", "*.docx"), 
            ("PowerPoint", "*.pptx")
        ]
    )
    
    if not caminho_arquivo:
        return # Execução silenciosamente encerrada se o usuário clicar em "Cancelar"
        
    try:
        # 2. Identifica o diretório (pasta) do arquivo selecionado
        diretorio_base = os.path.dirname(caminho_arquivo)
        nome_arquivo = os.path.basename(caminho_arquivo)
        extensao = nome_arquivo.lower()
        
        # 3. Busca obrigatoriamente o 'dados.xlsx' no mesmo diretório
        caminho_excel = os.path.join(diretorio_base, "dados.xlsx")
        
        if not os.path.exists(caminho_excel):
            messagebox.showwarning(
                "Aviso: Planilha não encontrada", 
                f"O arquivo 'dados.xlsx' não foi encontrado na pasta do relatório:\n\n{diretorio_base}\n\nPor favor, adicione o arquivo de dados e tente novamente."
            )
            return
            
        # 4. Lê as regras do Excel e define onde salvar o arquivo modificado
        regras = ler_regras_excel(caminho_excel)
        
        if not regras:
            messagebox.showinfo("Aviso", "A planilha 'dados.xlsx' foi encontrada, mas parece estar vazia ou sem regras válidas.")
            return

        caminho_saida = os.path.join(diretorio_base, f"Modificado_{nome_arquivo}")
        
        # 5. Processa de acordo com o formato
        if extensao.endswith('.docx'):
            processar_word(caminho_arquivo, caminho_saida, regras)
        elif extensao.endswith('.pptx'):
            processar_powerpoint(caminho_arquivo, caminho_saida, regras)
        else:
            messagebox.showerror("Erro", "Formato de arquivo não suportado. Escolha um .docx ou .pptx.")
            return
            
        # 6. Exibe mensagem de sucesso ao final
        messagebox.showinfo("Sucesso", f"Relatório modificado com sucesso!\n\nFoi gerado um novo arquivo na mesma pasta:\n{caminho_saida}")
        
    except PermissionError:
        messagebox.showerror(
            "Erro de Permissão", 
            "O arquivo original ou o arquivo 'dados.xlsx' está aberto em outro programa.\n\nFeche o Word, PowerPoint ou Excel e tente novamente."
        )
    except Exception as e:
        messagebox.showerror("Erro de Processamento", f"Ocorreu um erro durante a execução:\n\n{str(e)}")

if __name__ == "__main__":
    iniciar_automacao()